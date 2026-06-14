#!/usr/bin/env python3
"""Populate the Redrob idea-submission template with the Shortlist deck content.

Reads the blank template from the parent folder, fills each slide's guidance box with real,
explained content (bold lead + full sentence, not keyword fragments), draws a clean
single-row architecture pipeline on slide 7, and leaves the template's own closing slide 11
intact. Writes deck/Shortlist_Redrob.pptx. Local tooling only. Run: python deck/build_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "..", "..", "Idea Submission Template Redrob.pptx")
OUT = os.path.join(HERE, "Shortlist_Redrob.pptx")

INK = RGBColor(0x22, 0x26, 0x2B)
ACCENT = RGBColor(0x1F, 0x6F, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTE = RGBColor(0x55, 0x5B, 0x62)
ARROW = RGBColor(0x7A, 0x82, 0x8C)


def set_points(shape, items, size=11.0, lead_color=ACCENT):
    """Fill a shape with points. Each item is (lead, body): lead is a short bold phrase shown
    inline, body is a full explanatory sentence. Either may be None."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, (lead, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        p.line_spacing = 1.03
        if lead:
            rl = p.add_run(); rl.text = lead + "  "
            rl.font.size = Pt(size); rl.font.bold = True; rl.font.color.rgb = lead_color
        if body:
            rb = p.add_run(); rb.text = body
            rb.font.size = Pt(size); rb.font.color.rgb = INK


def guidance_shape(slide):
    cand = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    return max(cand, key=lambda s: len(s.text_frame.text)) if cand else None


def box(slide, x, y, w, h, text, fill, fg=WHITE, size=9.5, bold=True, line=ACCENT):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2); tf.margin_left = Pt(3); tf.margin_right = Pt(3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = fg
    return sp


def arrow(slide, shape, x, y, w, h):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = ARROW; sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


prs = Presentation(TEMPLATE)
S = prs.slides

# --- Slide 1: title ---
for sh in S[0].shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t.startswith("Team Name"):
        sh.text_frame.text = "Team Name :  Shortlist"
    elif t.startswith("Problem Statement"):
        sh.text_frame.text = ("Problem Statement :  Intelligent Candidate Discovery & Ranking. "
                              "From 100,000 candidates, surface the top 100 best-fit profiles for one "
                              "Senior AI Engineer role, ranked by genuine fit rather than keyword overlap.")
    elif t.startswith("Team Leader"):
        sh.text_frame.text = "Team Leader Name :  Nandish S  (solo)"

# --- Slide 2: Solution Overview & differentiation ---
set_points(guidance_shape(S[1]), [
    (None, "We rank candidates the way an experienced recruiter would: by what they have actually "
           "built, not by which keywords happen to appear on a profile."),
    ("Transparent core.", "Each final score is relevance_core x a behavioral-availability factor x a "
           "honeypot factor, so every rank breaks down into named numbers we can read and defend."),
    ("Semantic recall, capped.", "A TF-IDF JD-similarity term (weight 0.12) helps surface strong "
           "profiles written in plain language, but is weighted low on purpose so it can never "
           "override real career evidence."),
    ("Keyword-proof.", "Stuffing a profile with AI terms does nothing, because a skill earns credit "
           "only when the work history corroborates it."),
    ("Trap-aware and explainable.", "It separates genuine seniors from stuffers, look-alike "
           "'behavioral twins' and planted impossible profiles, and no model runs at scoring time, "
           "so there is nothing to hallucinate. It finishes in about 55 seconds on a CPU."),
], size=11.5)

# --- Slide 3: JD Understanding & Candidate Evaluation ---
set_points(guidance_shape(S[2]), [
    ("Requirements we extracted from the JD.", "5 to 9 years of experience (ideally 6 to 8); hands-on "
           "production work in ranking, search, recommendation or retrieval; strong Python and modern "
           "ML (PyTorch, embeddings, RAG, vector search); real offline evaluation (NDCG, A/B testing); "
           "a product-company rather than pure-services background; and being based in India or open to relocating."),
    ("Signals we weight most.", "Evidence of genuine information-retrieval work, skills we can actually "
           "trust, and whether the candidate sits inside the target experience band; availability "
           "signals then decide between otherwise-equal profiles."),
    ("How we go beyond keywords.", "A skill counts only when the career history backs it up; semantic "
           "similarity surfaces excellent candidates who never write buzzwords; and suspicious profiles "
           "are caught by internal contradictions, not by what they claim about themselves."),
], size=11.5)

# --- Slide 4: Ranking Methodology ---
set_points(guidance_shape(S[3]), [
    ("Retrieve.", "We stream and score all 100,000 candidates with no lossy pre-filter, so recall stays "
           "at 100% (the pool is small enough to score in full within the time budget)."),
    ("Score.", "relevance_core is a weighted sum of named signals: career evidence 0.30, semantic "
           "JD-similarity 0.12, trusted skills 0.18, experience band 0.12, Python and evaluation "
           "experience 0.12, location 0.12, and education 0.04."),
    ("Adjust.", "We multiply by a behavioral-availability factor (0.85 to 1.05) from recency, recruiter "
           "response, notice period and open-to-work status, and by a honeypot factor that pushes "
           "internally impossible profiles down."),
    ("Rank.", "We sort by the final score and break ties deterministically by candidate_id, then take "
           "the top 100. Because every signal is combined arithmetically, we can explain exactly why "
           "any candidate sits where they do."),
], size=11.5)

# --- Slide 5: Explainability & Data Validation ---
set_points(guidance_shape(S[4]), [
    ("Explanations that match the score.", "Each candidate's reasoning is generated from the exact "
           "feature values that produced their rank, so the explanation can never disagree with it."),
    ("No hallucination, by construction.", "No language model runs at scoring time (the network is off), "
           "and every clause must cite a real field such as title, years, trusted skills, product "
           "industry or notice period."),
    ("Honest about weaknesses.", "We surface concerns, not just strengths: low recruiter response, long "
           "inactivity, onsite-only-abroad or a long notice period are stated plainly."),
    ("Incomplete data handled fairly.", "Sentinel values (a -1 GitHub or offer-history score) are treated "
           "as 'missing', never as the worst case, so a candidate is not punished for an unlinked account."),
    ("Suspicious profiles gated.", "A conjunctive honeypot gate needs at least two independent "
           "impossibility signals before it acts, so genuine seniors are never wrongly demoted; the "
           "result is 0% honeypots in the top 100 against a 10% disqualification line."),
], size=10.5)

# --- Slide 6: End-to-End Workflow ---
set_points(guidance_shape(S[5]), [
    ("1. Offline precompute (untimed).", "build_artifacts.py streams the pool and builds a TF-IDF table, "
           "four JD seed vectors and a data-derived calibration anchor, saved as sim_model.json (152 KB, committed)."),
    ("2. Ranking (timed, under 5 minutes).", "rank.py streams every candidate, scores each one "
           "(relevance_core + behavioral + honeypot), sorts with the candidate_id tie-break, and writes "
           "the top-100 CSV with a reasoning string per candidate."),
    ("3. Validation.", "The official validator runs alongside our self-audit, which checks the honeypot "
           "rate, ID membership in the pool, output format and score distribution before any submission."),
    ("4. Live demo, same code.", "The hosted sandbox imports the very same scorer module, so the "
           "interactive demo can never drift from what the grader reproduces."),
], size=11.5)

# --- Slide 7: System Architecture (clean single-row pipeline) ---
PALE = RGBColor(0xEC, 0xF1, 0xF8)
y = 1.55
bw, gap, h = 1.66, 0.28, 0.95
xs = (10.0 - (5 * bw + 4 * gap)) / 2.0   # center the 5-box row with equal side margins
labels = ["candidates\n.jsonl(.gz)\n100k profiles",
          "Streaming reader\n+ feature\nextraction",
          "Scorer\nrelevance_core\n+ capped jd_similarity",
          "x behavioral\nx honeypot\nfactor",
          "Sort + candidate_id\ntie-break"]
fills = [PALE, PALE, ACCENT, ACCENT, PALE]
fgs = [INK, INK, WHITE, WHITE, INK]
centers = []
x = xs
for i, lab in enumerate(labels):
    box(S[6], x, y, bw, h, lab, fills[i], fg=fgs[i], size=9)
    centers.append(x + bw / 2)
    x += bw
    if i < len(labels) - 1:
        arrow(S[6], MSO_SHAPE.RIGHT_ARROW, x, y + h / 2 - 0.10, gap, 0.20)
        x += gap
# Output drops below the last box.
last_c = centers[-1]
arrow(S[6], MSO_SHAPE.DOWN_ARROW, last_c - 0.12, y + h + 0.05, 0.24, 0.32)
box(S[6], last_c - 1.0, y + h + 0.42, 2.0, 0.78, "submission.csv\n+ reasoning",
    RGBColor(0xDF, 0xF0, 0xE2), fg=INK, size=9, line=RGBColor(0x4C, 0xA1, 0x6B))
# Offline model feeds UP into the scorer (3rd box).
scorer_c = centers[2]
box(S[6], scorer_c - 1.15, y + h + 0.42, 2.3, 0.78, "build_artifacts.py\nsim_model.json (offline)",
    RGBColor(0xFB, 0xF1, 0xDC), fg=INK, size=9, line=RGBColor(0xC9, 0x9A, 0x2E))
arrow(S[6], MSO_SHAPE.UP_ARROW, scorer_c - 0.12, y + h + 0.05, 0.24, 0.32)
# Sandbox note (serving), bottom-left.
box(S[6], 0.20, y + h + 0.42, 2.0, 0.78, "Sandbox (Streamlit)\nimports same scorer",
    RGBColor(0xF3, 0xE9, 0xF7), fg=INK, size=9, line=RGBColor(0x9A, 0x6F, 0xB5))
# Wrapped caption, full width, below the diagram.
cap = S[6].shapes.add_textbox(Inches(0.20), Inches(4.35), Inches(9.4), Inches(0.7))
cap.text_frame.word_wrap = True
cp = cap.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
cr = cp.add_run()
cr.text = ("The ranking step uses only the Python standard library: no GPU, no network, no model "
           "download. The committed 152 KB similarity model is the only precomputed artifact.")
cr.font.size = Pt(10); cr.font.italic = True; cr.font.color.rgb = MUTE

# --- Slide 8: Results & Performance ---
set_points(guidance_shape(S[7]), [
    ("Strong, in-band head.", "The top 10 are all in-band IR/ML engineers at product companies (Search, "
           "Recommendation, NLP, ML and AI Research roles, 6 to 9 years); strong plain-language profiles "
           "surface even without buzzwords."),
    ("Stuffers do not survive.", "Of 5,588 keyword stuffers detected in the pool, none reach the top 100 "
           "(the best lands only at rank 628)."),
    ("Traps handled.", "0% honeypots in the top 100 (against a 10% disqualification line), and 100% recall "
           "on the hand-checked golden archetypes."),
    ("Comfortably within every constraint.", "About 55 seconds on 100k candidates (limit 5 minutes), well "
           "under 16 GB RAM, CPU-only, no network, and a ~150 KB model (limit 5 GB disk). It reproduces "
           "from a single command on a clean checkout."),
], size=11.5)

# --- Slide 9: Technologies Used ---
set_points(guidance_shape(S[8]), [
    ("Python standard library only (ranking step).", "json, gzip, csv, re and math. Zero ML dependencies "
           "means the timed step is trivially reproducible in the grader's Docker, fast, and free of "
           "version or download risk."),
    ("Custom stdlib TF-IDF (semantic layer).", "No scikit-learn and no model download. We chose this "
           "because the dataset's role descriptions are templated, so heavyweight embeddings add cost and "
           "reduce explainability without adding signal; a judge-panel comparison scored the transparent "
           "engine above an embeddings hybrid."),
    ("Streamlit.", "Powers the hosted sandbox demo only; it is never imported by the ranking step."),
    ("Git.", "Used for authentic, incremental version history that documents how the solution was built."),
], size=11.0)

# --- Slide 10: Submission Assets (guidance box is short, target it explicitly) ---
asset_box = next(s for s in S[9].shapes
                 if s.has_text_frame and "github video" in s.text_frame.text.lower())
set_points(asset_box, [
    ("GitHub repository.", "https://github.com/Nandish3010/redrob-candidate-ranker"),
    ("Hosted sandbox.", "https://redrob-candidate-ranker-ns.streamlit.app/"),
    ("Ranked output.", "outputs/submission.csv (top 100; columns candidate_id, rank, score, reasoning)."),
    ("Single reproduce command.", "python rank.py --candidates ./data/candidates.jsonl --out ./outputs/submission.csv"),
], size=12.0)

# --- Slide 11: leave the template's own closing design intact (no overlay added). ---

prs.save(OUT)
print("Wrote", OUT)
